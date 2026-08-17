# -*- coding: utf-8 -*-
"""Sinh slide trinh bay 24 trang + kich ban noi (ghi vao phan notes cua tung slide)."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = pathlib.Path(__file__).resolve().parent.parent
OUT = ROOT / "slide" / "slide-trinh-bay.pptx"
IMG = ROOT / "model" / "hinh-xuat"
FIG11 = ROOT / "docs" / "kien-truc-quy-trinh" / "hinh-1-1-kien-truc-quy-trinh.png"

NAVY, GREY, ACCENT = RGBColor(0x1B, 0x35, 0x5E), RGBColor(0x55, 0x55, 0x55), RGBColor(0xC0, 0x7A, 0x10)
FONT = "Segoe UI"


def txt(tf, lines, size=18, color=GREY, bullet=True):
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        lvl = 0
        if isinstance(ln, tuple):
            ln, lvl = ln
        p.level = lvl
        r = p.add_run()
        r.text = ("• " if bullet and lvl == 0 else ("– " if bullet else "")) + ln
        r.font.size, r.font.name, r.font.color.rgb = Pt(size - lvl * 2), FONT, color
        p.space_after = Pt(8)


def slide(prs, title, body=None, note="", img=None, size=18):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size, r.font.bold, r.font.name, r.font.color.rgb = Pt(30), True, FONT, NAVY
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.25), Inches(12.1), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    if img and pathlib.Path(img).exists():
        s.shapes.add_picture(str(img), Inches(0.6), Inches(1.7), width=Inches(12.1))
        if body:
            b = s.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.9))
            txt(b.text_frame, body, 15)
    elif body:
        b = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
        txt(b.text_frame, body, size)
    s.notes_slide.notes_text_frame.text = note
    return s


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for y, t, sz, b, c in [(2.2, "PHÂN TÍCH QUY TRÌNH NGHIỆP VỤ", 40, True, NAVY),
                           (3.0, "Chuỗi bán lẻ thegioididong.com và TopZone", 26, False, GREY),
                           (4.2, "Nhóm 4 · Hệ thống Quản trị Qui trình Nghiệp vụ", 18, False, GREY),
                           (4.7, "GVHD: ThS. Hà Lê Hoài Trung", 18, False, GREY),
                           (5.4, "Nguyễn Ngọc Danh · Nguyễn Thị Hồng Phúc · "
                                 "Nguyễn Thanh Phúc · Mai Hoàng Hưng", 15, False, GREY)]:
        tb = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.8))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name, r.font.color.rgb = Pt(sz), b, FONT, c
    s.notes_slide.notes_text_frame.text = (
        "Chào thầy và các bạn. Nhóm em trình bày đề tài phân tích quy trình nghiệp vụ "
        "chuỗi bán lẻ Thế Giới Di Động và TopZone. Bài gồm bốn phần: tổng quan và kiến "
        "trúc quy trình, cơ sở lý thuyết, mô hình hóa, và phân tích cải tiến.")


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    title_slide(prs)                                                        # 1

    slide(prs, "Nội dung trình bày", [
        "Phần 1 — Tổng quan doanh nghiệp và kiến trúc quy trình",
        "Phần 2 — Cơ sở lý thuyết và phương pháp",
        "Phần 3 — Mô hình hóa quy trình bằng BPMN",
        "Phần 4 — Phân tích, phát hiện và đề xuất cải tiến"],
        "Bài chia bốn phần, tổng khoảng 15 phút. Em sẽ dành nhiều thời gian nhất cho "
        "phần 3 và 4 vì đó là phần có kết quả phân tích.")                  # 2

    slide(prs, "1.1 Doanh nghiệp và phạm vi", [
        "MWG vận hành hai chuỗi: thegioididong.com (phổ thông) và TopZone (ủy quyền Apple)",
        "Dùng chung nền tảng ERP, kho và hệ thống bảo hành",
        "Phạm vi sản phẩm: điện thoại, laptop, máy tính bảng, phụ kiện",
        "Loại khỏi phạm vi: đồng hồ, hàng gia dụng, hàng tiêu dùng, dịch vụ SIM"],
        "Hai chuỗi dùng chung nền tảng vận hành nên quy trình lõi gộp được, chỉ tách khi "
        "có khác biệt thật sự — ví dụ khâu bảo hành của TopZone đi theo kênh Apple.")   # 3

    slide(prs, "1.2 Giới hạn — nói trước để mọi số liệu sau đều truy được nguồn", [
        "Nhóm là người quan sát bên ngoài, không truy cập hệ thống nội bộ MWG",
        "Số liệu thời gian đến từ bấm giờ tại cửa hàng (ghi rõ cỡ mẫu n)",
        "Số không đo được ghi rõ \u201c(ước lượng)\u201d và đưa vào bảng giả định",
        "Quy trình mô tả là mô hình suy luận, không phải quy trình MWG ban hành"],
        "Đây là slide em muốn nói kỹ. Nhóm không có số liệu nội bộ, nên thay vì bịa ra "
        "con số đẹp, nhóm ghi rõ cái gì đo được và cái gì không. Mọi con số trong bài "
        "đều truy được về bảng bấm giờ hoặc bảng giả định.")                # 4

    slide(prs, "1.3 Kiến trúc 12 quy trình theo ba lớp", [
        "6 quy trình được mô hình hóa BPMN: M2, M3, C3, C4, S1, S4"],
        "Đây là Hình 1.1 trong báo cáo. Ba dải là ba lớp. Mũi tên đỏ là vòng phản hồi từ "
        "bảo hành ngược lên quản lý nhà cung cấp — em sẽ quay lại chỗ này ở phần 4.",
        img=FIG11)                                                          # 5

    slide(prs, "1.3 Vì sao vòng phản hồi C4 → M2 quan trọng", [
        "Dữ liệu lỗi thực tế phát sinh ở khâu bảo hành (C4)",
        "Chỉ có ích cho việc chọn nhà cung cấp (M2) nếu quay về kịp thời",
        "Trễ ở vòng này = tiếp tục nhập model tỷ lệ lỗi cao",
        "Chi phí bảo hành cộng dồn qua nhiều chu kỳ nhập hàng",
        "Ghi nhận thành phát hiện IR-13 trong Issue Register"],
        "Đây là luồng duy nhất đi ngược từ lớp cốt lõi lên lớp quản lý, và là phát hiện "
        "mà nhóm thấy đáng giá nhất về mặt kiến trúc.")                     # 6

    for t, b, n in [
        ("2.1 Khái niệm quy trình nghiệp vụ và BPM lifecycle", None,
         "Phần lý thuyết do bạn Hồng Phúc trình bày."),
        ("2.2 Ký hiệu BPMN 2.0 dùng trong đề tài", None,
         "Bảng ký hiệu — Hồng Phúc trình bày."),
        ("2.3 Ba nhóm lãng phí Move / Hold / Overdo", None,
         "Khung phân loại lãng phí — Hồng Phúc trình bày.")]:
        slide(prs, t, b or ["[Nội dung do Hồng Phúc bổ sung]"], n)          # 7,8,9

    slide(prs, "2.4 Phương pháp thu thập dữ liệu", [
        "Quan sát trực tiếp tại cửa hàng TGDĐ",
        "Bấm giờ 6–10 lượt giao dịch bán, bấm riêng khâu bảo hành",
        "Phỏng vấn nhân viên — 6 câu ngắn, không ghi âm nếu không được đồng ý",
        "Chụp biểu mẫu và chính sách, che thông tin cá nhân",
        "[Số liệu cụ thể do Hưng bổ sung sau buổi khảo sát]"],
        "Bạn Hưng phụ trách buổi khảo sát thực địa. Đây là phần khác biệt lớn nhất của "
        "nhóm so với việc chỉ đọc tài liệu trên mạng.")                     # 10

    slide(prs, "2.5 Review chéo mô hình — bằng chứng từ pull request", [
        "Mỗi mô hình: một người vẽ, một người khác duyệt qua pull request",
        "Comment trong PR dùng luôn làm biên bản review chéo",
        "Tự kiểm bằng script trên file .bpmn, không đếm tay",
        "M2: 9 gateway · C3: 10 gateway · C4: 11 gateway — đều vượt mốc 7"],
        "Nhóm không merge im lặng. Comment trong pull request chính là biên bản review, "
        "một việc ra hai sản phẩm.")                                        # 11
    return prs


def phan_3_4(prs):
    slide(prs, "3.1 Ba quy trình cốt lõi — hồ sơ đầy đủ 12 mục", [
        "C1 Bán tại cửa hàng — 14 bước, 9 điểm quyết định, 6 ngoại lệ",
        "C3 Bán trả góp — 15 bước, 10 điểm quyết định, 7 ngoại lệ",
        "C4 Bảo hành đổi trả — 16 bước, 11 điểm quyết định, 7 ngoại lệ",
        "Template thống nhất: mục đích, phạm vi, actor, I/O, bước, gateway, ngoại lệ, "
        "quy tắc, chỉ số, hệ thống, điểm nghẽn, nguồn"],
        "Ba hồ sơ dùng chung một template 12 mục, nên Chương 3 dàn trang đồng nhất và "
        "bảng nào cũng so sánh được với nhau.")                             # 12

    slide(prs, "3.4 BPMN M2 — Quản lý nhà cung cấp (9 gateway)", [
        "Lưu ý: dựng khi hồ sơ M2 chưa hoàn tất, cần đối chiếu trước 30/08"],
        "Mô hình M2 có vòng phê duyệt nhiều cấp theo giá trị hợp đồng, và vòng đánh giá "
        "nhà cung cấp định kỳ có dùng dữ liệu lỗi từ C4.",
        img=IMG / "M2-quan-ly-nha-cung-cap.png")                            # 13

    slide(prs, "3.4 BPMN C3 — Bán trả góp (10 gateway)", [
        "G4 ba nhánh: duyệt / từ chối / yêu cầu bổ sung"],
        "Điểm đáng chú ý ở C3 là actor công ty tài chính nằm ngoài doanh nghiệp nhưng "
        "lại là người ra quyết định. Toàn bộ quy trình dừng chờ ở đó.",
        img=IMG / "C3-ban-tra-gop.png")                                     # 14

    slide(prs, "3.4 BPMN C4 — Bảo hành, đổi trả (11 gateway)", [
        "Trung tâm bảo hành vẽ thành lane riêng — actor bên ngoài ra quyết định cuối"],
        "C4 là mô hình phức tạp nhất với 11 gateway. Có bốn kết thúc khác nhau, trong đó "
        "hai kết thúc là từ chối bảo hành.",
        img=IMG / "C4-bao-hanh-doi-tra.png")                                # 15

    for t, n in [("4.1 Bảng giả định cho số liệu chưa xác minh", "Hưng trình bày."),
                 ("4.2 Cycle time và hiệu suất thời gian CTE", "Hưng trình bày."),
                 ("4.3 Phân tích giá trị VA / BVA / NVA cho C3 và C4", "Hồng Phúc trình bày."),
                 ("4.4 Phân nhóm lãng phí Move / Hold / Overdo", "Hồng Phúc trình bày."),
                 ("4.5 Biểu đồ xương cá cho điểm nghẽn chính", "Hồng Phúc trình bày.")]:
        slide(prs, t, ["[Nội dung bổ sung sau buổi khảo sát 23/08]"], n)    # 16–20

    slide(prs, "4.6 Issue Register — 13 phát hiện", [
        "Hold — chờ đợi: 6 phát hiện (IR-03, 04, 05, 09, 10, 13)",
        "Overdo — làm thừa, làm lại: 6 phát hiện (IR-01, 02, 06, 07, 11, 12)",
        "Move — di chuyển: 2 phát hiện (IR-08, IR-12)",
        "Phân bố theo quy trình: C1 bốn · C3 bốn · C4 năm"],
        "Mười ba phát hiện gom từ mục điểm nghẽn của ba hồ sơ. Nhóm Hold chiếm gần một "
        "nửa — đó là hướng chính của phần đề xuất.")                        # 21

    slide(prs, "4.6 Hai loại chờ đợi cần hai hướng xử lý khác nhau", [
        "Chờ actor bên ngoài — IR-05 công ty tài chính, IR-10 trung tâm bảo hành",
        ("Cửa hàng không kiểm soát được → giảm ảnh hưởng, không xóa được khoảng chờ", 1),
        "Chờ nguồn lực nội bộ — IR-03, IR-04, IR-09",
        ("Cửa hàng kiểm soát được → bố trí lại nhân sự và điểm phục vụ", 1),
        "Xếp ưu tiên theo hai trục: mức ảnh hưởng × khả năng can thiệp"],
        "Đây là điểm em muốn nhấn: hai phát hiện nặng nhất về thời gian lại là hai phát "
        "hiện cửa hàng không tự xử lý được. Đề xuất phải nhắm vào giảm ảnh hưởng của "
        "khoảng chờ chứ không phải xóa nó.")                                # 22

    slide(prs, "4.7 Đề xuất cải tiến", [
        "[Viết sau khi có số liệu định lượng mục 4.2 — hạn 02/09]"],
        "Phần đề xuất chờ số liệu khảo sát để xếp ưu tiên có căn cứ.")      # 23

    slide(prs, "Kết luận", [
        "Phân rã 12 quy trình theo ba lớp, mô hình hóa 6 quy trình bằng BPMN",
        "Ba mô hình của nhóm trưởng: 9, 10 và 11 gateway — đều vượt mốc 7",
        "13 phát hiện được gom thành Issue Register có mã truy vết",
        "Mọi số liệu chưa xác minh đều đánh dấu rõ, không bịa số nội bộ",
        "Hạn chế: quy trình dựng từ quan sát bên ngoài, cần đối chiếu thực địa"],
        "Em xin kết thúc phần trình bày. Nhóm em cảm ơn thầy và các bạn đã lắng nghe.") # 24
    return prs


if __name__ == "__main__":
    p = phan_3_4(build())
    p.save(OUT)
    n_note = sum(1 for s in p.slides if s.notes_slide.notes_text_frame.text.strip())
    print("Da ghi %s | %d slide | %d slide co kich ban noi"
          % (OUT.name, len(p.slides._sldIdLst), n_note))
